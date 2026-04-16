"""
Ingest documents from context/ and workspace/ folders.
Filters out binary assets (images, fonts, compiled files, etc.).
Supports: code, text, pptx, docx, pdf.

Each chunk is tagged with a `doc_level` metadata for hierarchical RAG:
  - L0: Architecture overview (synthesize.py output)
  - L1: Layer overviews (synthesize.py output)
  - L2: Module docs (synthesize.py output)
  - L3: Codex scan docs (codex_*.md)
  - code: Raw source code from workspace
  - context: Manual docs (architecture/, code-samples/)
  - report: Agent reports

Features:
  - Semantic breadcrumbs for better embedding context
  - Rich metadata (block, module, content_type)
  - CRLF normalization
  - Incremental ingestion with file hashing
"""

import hashlib
import json
import logging
import os
import re
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# -- Binary / asset extensions to SKIP ---

SKIP_DIRS = {
    "node_modules", "__pycache__", ".git", ".svn", ".hg",
    "dist", "build", ".next", ".nuxt", "target",
    ".venv", "venv", "env", ".env",
    ".vectordb", ".idea", ".vscode",
    "vendor", "bower_components",
}

MAX_FILE_SIZE = 1_000_000  # 1MB


# -- Parsers ---

def _read_text(path: Path) -> str:
    try:
        content = path.read_text(encoding="utf-8", errors="replace")
        if content.count("\ufffd") > len(content) * 0.1:
            logger.debug(f"Skipping likely binary file: {path}")
            return ""
        return content
    except Exception as e:
        logger.warning(f"Cannot read {path}: {e}")
        return ""


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_texts.append(shape.text_frame.text)
            if slide_texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except Exception as e:
        logger.warning(f"Cannot parse PPTX {path}: {e}")
        return ""


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.warning(f"Cannot parse DOCX {path}: {e}")
        return ""


def _read_pdf(path: Path) -> str:
    try:
        import fitz
        doc = fitz.open(str(path))
        parts = []
        for page in doc:
            parts.append(page.get_text())
        doc.close()
        return "\n\n".join(parts)
    except Exception as e:
        logger.warning(f"Cannot parse PDF {path}: {e}")
        return ""


PARSERS = {
    ".pptx": _read_pptx,
    ".docx": _read_docx,
    ".pdf": _read_pdf,
}


# -- Doc level detection ---

def detect_doc_level(relative_path: str, directory_label: str = "") -> str:
    """
    Detect the documentation level of a file for hierarchical RAG.
    
    Args:
        relative_path: Path relative to the ingestion directory
        directory_label: Label passed by the caller ("context", "workspace", "reports")
    
    Returns:
        One of: "L0", "L1", "L2", "L3", "code", "context", "report"
    """
    rel_lower = relative_path.lower()
    name = Path(relative_path).name.lower()

    # Reports
    if directory_label == "reports":
        return "report"

    # Workspace = raw code
    if directory_label == "workspace":
        return "code"

    # Context directory: detect synthesis levels
    if directory_label == "context":
        # Synthesis outputs
        if "synthesis/" in rel_lower or "synthesis\\" in rel_lower:
            m = re.match(r"^(l\d+)_", name)
            if m:
                return m.group(1).upper()
        # Codex scan docs
        if name.startswith("codex_"):
            return "L3"

        # Manual context docs
        return "context"

    # Fallback
    return "context"


# -- Incremental ingestion state ---

def _load_hashes(ingest_dir: Path) -> dict:
    """Load file hashes from sidecar JSON file."""
    hash_file = ingest_dir / ".ingest_hashes.json"
    if hash_file.exists():
        try:
            return json.loads(hash_file.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, OSError):
            pass
    return {}


def _save_hashes(ingest_dir: Path, hashes: dict) -> None:
    """Save file hashes to sidecar JSON file."""
    hash_file = ingest_dir / ".ingest_hashes.json"
    hash_file.write_text(
        json.dumps(hashes, indent=2, ensure_ascii=False),
        encoding="utf-8"
    )


def _compute_file_hash(content: str) -> str:
    """Compute MD5 hash of normalized content."""
    normalized = content.replace("\r\n", "\n").replace("\r", "\n")
    return hashlib.md5(normalized.encode("utf-8")).hexdigest()


# -- Chunking with semantic breadcrumbs ---

def chunk_text_with_breadcrumb(
    text: str,
    chunk_size: int = 1500,
    overlap: int = 200,
    source: str = "",
    doc_level: str = "context",
    block: str = "",
    module_name: str = "",
    content_type: str = "",
) -> list[dict]:
    """
    Split text into chunks with metadata including semantic breadcrumbs.
    
    Each chunk dict contains: text, source, chunk_index, doc_level, and rich metadata.
    
    Args:
        text: The text content to chunk
        chunk_size: Target chunk size in characters
        overlap: Overlap between chunks
        source: Source file path
        doc_level: Documentation level (L0, L1, L2, L3, code, context, report)
        block: Architectural block (backend, frontend, database, etc.)
        module_name: Module name
        content_type: Type of content (code, codex_doc, synthesis, config, test)
    
    Returns:
        List of chunk dicts with rich metadata
    """
    if not text.strip():
        return []

    # Normalize line endings
    normalized_text = text.replace("\r\n", "\n").replace("\r", "\n")
    
    # Add semantic breadcrumb context line
    breadcrumb = (
        f"Context: Block={block} | Module={module_name} | Level={doc_level} | Type={content_type}\n\n"
    )
    normalized_text = breadcrumb + normalized_text

    chunks = []
    start = 0
    idx = 0
    while start < len(normalized_text):
        end = start + chunk_size
        chunk = normalized_text[start:end]

        if end < len(normalized_text):
            last_nl = chunk.rfind("\n")
            if last_nl > chunk_size // 2:
                end = start + last_nl + 1
                chunk = normalized_text[start:end]

        clean_chunk = chunk.strip().replace('\x00', '').replace('\0', '')
        if clean_chunk:
            chunks.append({
                "text": clean_chunk,
                "source": str(source) if source is not None else "unknown",
                "chunk_index": int(idx),
                "doc_level": str(doc_level) if doc_level is not None else "unknown",
                "block": block,
                "module": module_name,
                "content_type": content_type,
            })
        start = end - overlap
        idx += 1

    return chunks


# -- Ingestion pipeline ---

def _should_skip_dir(dirname: str, skip_dirs: set[str] = SKIP_DIRS) -> bool:
    return dirname in skip_dirs or dirname.startswith(".")


def _should_skip_file(path: Path, allowed_extensions: set[str], max_file_size: int = MAX_FILE_SIZE) -> bool:
    """Skip files that don't match our criteria.

    - Must be in allowed_extensions (if specified)
    - Must not exceed max_file_size
    - Minified files are always skipped
    """
    suffix = path.suffix.lower()

    # Skip minified files (always)
    if path.name.endswith((".min.js", ".min.css", ".bundle.js", ".chunk.js")):
        return True

    # Skip by size
    try:
        if path.stat().st_size > max_file_size:
            return True
    except OSError:
        return True

    # Positive list: only index if extension is in allowed_extensions
    if allowed_extensions and suffix not in allowed_extensions:
        return True

    return False


def ingest_directory(
    directory: str | Path,
    extensions: Optional[list[str]] = None,
    chunk_size: int = 1500,
    chunk_overlap: int = 200,
    label: str = "",
    skip_dirs: Optional[set[str]] = None,
    max_file_size: Optional[int] = None,
    force: bool = False,
    ingest_dir: Optional[Path] = None,
) -> list[dict]:
    """
    Recursively read all matching files, parse, and chunk.

    Args:
        directory: Directory to scan
        extensions: Allowed file extensions (None = all)
        chunk_size: Chunk size in chars
        chunk_overlap: Overlap between chunks
        label: Directory label for doc_level detection ("context", "workspace", "reports")
        skip_dirs: Directories to skip
        max_file_size: Maximum file size
        force: Force re-ingestion even if file hasn't changed
        ingest_dir: Directory for storing hash state (defaults to directory)

    Returns:
        List of {text, source, chunk_index, doc_level, block, module, content_type} dicts.
    """
    directory = Path(directory)
    if not directory.exists():
        logger.warning(f"Directory {directory} does not exist")
        return []

    _skip_dirs = skip_dirs if skip_dirs is not None else SKIP_DIRS
    _max_file_size = max_file_size if max_file_size is not None else MAX_FILE_SIZE
    allowed = set(extensions) if extensions else set()
    all_chunks = []
    skipped = 0
    processed = 0

    # Determine ingest state directory
    if ingest_dir is None:
        ingest_dir = directory
    
    state = _load_hashes(ingest_dir)

    all_paths = []
    for dirpath, dirnames, filenames in os.walk(directory, followlinks=True):
        dirnames[:] = [d for d in dirnames if d not in _skip_dirs and not d.startswith(".")]
        for fname in filenames:
            all_paths.append(Path(dirpath) / fname)
    
    for path in sorted(all_paths):
        if not path.is_file():
            continue
        if any(_should_skip_dir(part, _skip_dirs) for part in path.relative_to(directory).parts[:-1]):
            continue
        if _should_skip_file(path, allowed, _max_file_size):
            skipped += 1
            continue

        # Check if file has changed
        try:
            current_content = path.read_text(encoding="utf-8", errors="replace")
            current_hash = _compute_file_hash(current_content)
            file_key = str(path.relative_to(ingest_dir))
            
            if not force and state.get(file_key) == current_hash:
                logger.debug(f"Skipping unchanged file: {path}")
                continue
        except Exception as e:
            logger.warning(f"Cannot read {path} for hashing: {e}")
            continue

        logger.info(f"Ingesting: {path}")
        processed += 1

        parser = PARSERS.get(path.suffix.lower(), _read_text)
        text = parser(path)
        if not text.strip():
            continue

        relative = path.relative_to(directory) if directory in path.parents or directory == path.parent else path.name
        relative_str = str(relative)
        header = f"[File: {relative_str}]\n"
        text_with_header = header + text

        doc_level = detect_doc_level(relative_str, label)
        
        # Extract block and module from path for semantic breadcrumbs
        block = "other"
        module_name = ""
        content_type = ""
        
        # Determine block based on directory structure
        rel_path = str(relative)
        if "backend" in rel_path or rel_path.endswith(".java") or rel_path.endswith(".py"):
            block = "backend"
        elif "frontend" in rel_path or rel_path.endswith(".js") or rel_path.endswith(".ts"):
            block = "frontend"
        elif "database" in rel_path or rel_path.endswith(".sql"):
            block = "database"
        elif "test" in rel_path or "spec" in rel_path:
            block = "tests"
        elif "infra" in rel_path or "deploy" in rel_path:
            block = "infrastructure"
        
        # Determine content type based on file
        if label == "workspace":
            content_type = "code"
        elif "codex_" in relative_str:
            content_type = "codex_doc"
        elif "synthesis" in rel_path:
            content_type = "synthesis"
        elif "config" in rel_path or rel_path.endswith((".yaml", ".yml", ".json")):
            content_type = "config"
        elif "test" in rel_path or "spec" in rel_path:
            content_type = "test"
        else:
            content_type = "context"
        
        # Extract module name from path
        if len(relative.parts) > 0:
            module_name = relative.parts[0]
        
        chunks = chunk_text_with_breadcrumb(
            text_with_header,
            chunk_size=chunk_size,
            overlap=chunk_overlap,
            source=relative_str,
            doc_level=doc_level,
            block=block,
            module_name=module_name,
            content_type=content_type,
        )
        all_chunks.extend(chunks)
        
        # Update hash state
        state[file_key] = current_hash

    # Save updated hash state
    if processed > 0:
        _save_hashes(ingest_dir, state)

    # Log level distribution
    level_counts = {}
    for c in all_chunks:
        lvl = c.get("doc_level", "?")
        level_counts[lvl] = level_counts.get(lvl, 0) + 1

    logger.info(
        f"Ingested {len(all_chunks)} chunks from {directory} "
        f"(processed {processed} files, skipped {skipped} binary/asset files) "
        f"levels: {level_counts}"
    )
    return all_chunks
